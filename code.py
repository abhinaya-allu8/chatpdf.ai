import streamlit as st
import os
import fitz  # PyMuPDF
from pptx import Presentation  # For handling PPTX files
import docx  # For handling DOCX files
from langchain_core.messages import HumanMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
import re

# Define the correct passkey
correct_passkey = "AIChat"  # Replace with your actual passkey

# Function to read a PDF file and extract its text content
def read_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# Function to read a PPTX file and extract its text content
def read_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to read a D
